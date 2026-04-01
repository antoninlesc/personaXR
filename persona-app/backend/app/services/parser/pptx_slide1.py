from pptx import Presentation
from .utils import collect_all_shapes, get_shape_text, get_shapes_in_direction

def parse_pptx(file_path: str) -> dict:
    prs = Presentation(file_path)
    slide = prs.slides[0]

    # Quick map by name/id for targeted extraction if standard template
    shapes_by_name = {}
    shapes_by_id = {}
    
    # Flatten all shapes (including groups)
    all_shapes = collect_all_shapes(slide.shapes)
    bio_shape = None
    inquietudes_shape = None
    besoins_shape = None
    insight_shape = None

    for shape in all_shapes:
        shapes_by_name[shape.name] = shape
        shapes_by_id[shape.shape_id] = shape

        if hasattr(shape, "text_frame") and shape.text_frame:
            if "bio" in shape.text_frame.text.lower() and not bio_shape:
                bio_shape = shape
            
            if "inquiétudes" in shape.text_frame.text.lower() and not inquietudes_shape:
                inquietudes_shape = shape
            
            if "besoins" in shape.text_frame.text.lower() and not besoins_shape:
                besoins_shape = shape
            
            if "insight" in shape.text_frame.text.lower() and not insight_shape:
                insight_shape = shape


    data = {
        "slide": 1,
        "nom": "",
        "metier": "",
        "localisation": "",
        "caracteristiques": "",
        "phrase_cle": "",
        "maturite_digitale": 0,
        "bio": "",
        "insights": [],
        "inquietudes": "", 
        "besoins": "",
        "meta": {"parser_version": "0.0.2"}
    }
    
    # 1. Nom (TextBox 9)
    # The name is unique in this template, so this is safe.
    t9 = next((s for s in all_shapes if s.name == "TextBox 9"), None)
    if t9:
        data["nom"] = get_shape_text(t9)

    # 2. Métier, Localisation (TextBox 83)
    t83 = next((s for s in all_shapes if s.name == "TextBox 83"), None)
    if t83:
        text = get_shape_text(t83)
        lines = [l.strip() for l in text.split('\n') if l.strip()]
        if len(lines) >= 1:
            # First line is Age usually
            pass
        if len(lines) >= 2:
            data["localisation"] = lines[1]
        if len(lines) >= 3:
            data["metier"] = lines[2]

    # 3. Caractéristiques (Rounded Rectangle 72, 74)
    r72 = next((s for s in all_shapes if s.name == "Rounded Rectangle 72"), None)
    r74 = next((s for s in all_shapes if s.name == "Rounded Rectangle 74"), None)
    
    chars = []
    if r72: chars.append(get_shape_text(r72))
    if r74: chars.append(get_shape_text(r74))
    data["caracteristiques"] = ", ".join([c for c in chars if c])

    # 4. Phrase clé (Rectangle 5) or similar
    # Sometimes quotes are in "TextBox 40" (the symbol) but the text is in "Rectangle 5".
    # Analysis showed Rectangle 5 has: "Je ne comprends pas..."
    r5 = next((s for s in all_shapes if s.name == "Rectangle 5"), None)
    if r5:
        data["phrase_cle"] = get_shape_text(r5)

    bio_text = get_shapes_in_direction(bio_shape, all_shapes) if bio_shape else []
    data["bio"] = "\n".join(bio_text)
    data["inquietudes"] = "\n".join(get_shapes_in_direction(inquietudes_shape, all_shapes))
    data["besoins"] = "\n".join(get_shapes_in_direction(besoins_shape, all_shapes))
    data["insights"] = get_shapes_in_direction(insight_shape, all_shapes, max_count=2)

    return data