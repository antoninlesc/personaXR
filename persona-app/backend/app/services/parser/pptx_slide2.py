from pptx import Presentation
from .utils import collect_all_shapes, get_shapes_in_direction, get_shapes, get_shapes_inside, is_below, return_position

def parse_slide2(file_path):
    prs = Presentation(file_path)
    
    # Vérification qu'il y a bien une slide 2
    if len(prs.slides) < 2:
        return {"error": "No slide 2 found"}
        
    slide = prs.slides[1] # Index 1 = Slide 2

    #all shapes and the groups are extended to a single list for easier processing
    all_shapes = collect_all_shapes(slide.shapes)

    #all shapes but the groups aren't extended
    every_shapes = get_shapes(slide)

    # Structure de données cible : Une liste d'étapes
    data = {
        "steps": [],
        "emotion_gen": []
    }

    # 1. Trouver les En-têtes de colonnes (Les étapes)
    # Heuristique : Ce sont souvent des formes situées en haut de la page, alignées horizontalement
    # On cherche les formes dans le quart supérieur
    potential_headers = []
    row_headers = []

    for shape in all_shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            if "etape" in shape.text_frame.text.lower() and potential_headers == []:
                text = get_shapes_in_direction(shape, all_shapes, max_count=0, direction="below", mode="shape")
                potential_headers = get_shapes_in_direction(shape, all_shapes, direction="right", max_count=0, mode="shape")
                row_headers.append(shape)
                row_headers.extend(text)
                break

    top_pos_emotion = return_position(row_headers[2])["bottom"] - 0.5
    bot_pos_emotion = return_position(row_headers[4])["top"] + 0.5
    actions = get_shapes_in_direction(row_headers[1], all_shapes, direction="right", max_count=0, mode="shape")
    challenges = get_shapes_in_direction(row_headers[4], all_shapes, direction="right", max_count=0, mode="shape")

    for shape in all_shapes:
        pos = return_position(shape)
        if pos["top"] >= top_pos_emotion and pos["bottom"] <= bot_pos_emotion:
            if hasattr(shape, "text") and shape.text.strip() and not shape.text.strip().isdigit() and "emotion" not in shape.text.strip().lower():
                data["emotion_gen"].append(shape.text.strip().replace("\n", " "))

    for header in potential_headers:
        step_data = {
            "name": header.text_frame.text.strip() if hasattr(header, "text_frame") and header.text_frame else header.name,
            "actions": [],
            "canaux": "",
            "emotions": "", 
            "challenges": []
        }

        for s in actions:
            insides = get_shapes_inside(s, every_shapes, mode='shape')
            inside = {
                s.name: []
            }
            for ins in insides:
                if is_below(header, ins):
                    if hasattr(ins, "shapes"):
                        for sub in ins.shapes:
                            if hasattr(sub, "text") and sub.text.strip() and not sub.text.strip().isdigit():
                                inside[s.name].append(sub.text.strip())
                    else :
                        inside[s.name].append(getattr(ins, "text").strip())
            step_data["actions"].extend(inside[s.name])

        for chal in challenges:
            if is_below(header, chal):
                if hasattr(chal, "text") and chal.text.strip():
                    step_data["challenges"].extend(chal.text.strip().split("\n")) # On prend la première ligne pour éviter les listes à puces

        data["steps"].append(step_data)


    return data