from pydoc import text
from pptx.util import Length

def get_shape_text(shape):
    if hasattr(shape, "text_frame") and shape.text_frame:
        return shape.text_frame.text.strip()
    return ""

def get_shapes(slide):
    shapes = []
    for shape in slide.shapes:
        shapes.append(shape)
    return shapes

def collect_all_shapes(shapes):
    """Recursively collect all shapes, including those in groups"""
    all_shapes = []
    for shape in shapes:
        if shape.shape_type == 6:
            all_shapes.extend(collect_all_shapes(shape.shapes))
        else:
            all_shapes.append(shape)
    return all_shapes

def get_shapes_in_direction(ref_shape, all_shapes, direction="below", max_count=1, mode="text"):
    """
    Trouve les formes situées dans une direction spécifique par rapport à ref_shape.
    Utilise une logique de 'Ray Casting' (Laser) depuis le centre de la ref_shape.
    """
    if not ref_shape:
        return []

    # Coordonnées de référence
    ref_center_x = ref_shape.left + (ref_shape.width / 2)
    ref_center_y = ref_shape.top + (ref_shape.height / 2)
    
    ref_left = ref_shape.left
    ref_right = ref_shape.left + ref_shape.width
    ref_top = ref_shape.top
    ref_bottom = ref_shape.top + ref_shape.height

    candidates = []

    for shape in all_shapes:
        # On ignore la forme elle-même
        if shape.shape_id == ref_shape.shape_id:
            continue
        
        # Coordonnées de la forme candidate
        s_left = shape.left
        s_right = shape.left + shape.width
        s_top = shape.top
        s_bottom = shape.top + shape.height

        match_found = False
        dist = float('inf')

        # --- LOGIQUE DIRECTIONNELLE ---
        
        if direction == "below":
            # 1. Est-ce en dessous ? (Top du candidat >= Bas de la ref)
            if s_top >= (ref_bottom - 5000): # Petite tolérance
                # 2. Est-ce aligné horizontalement ? (Le centre X de ref traverse le candidat)
                if s_left <= ref_center_x <= s_right:
                    match_found = True
                    dist = s_top - ref_bottom

        elif direction == "above":
            # 1. Est-ce au-dessus ? (Bas du candidat <= Haut de la ref)
            if s_bottom <= (ref_top + 5000):
                # 2. Est-ce aligné horizontalement ?
                if s_left <= ref_center_x <= s_right:
                    match_found = True
                    dist = ref_top - s_bottom

        elif direction == "right":
            # 1. Est-ce à droite ? (Gauche du candidat >= Droite de la ref)
            if s_left >= (ref_right - 5000):
                # 2. Est-ce aligné verticalement ? (Le centre Y de ref traverse le candidat)
                if s_top <= ref_center_y <= s_bottom:
                    match_found = True
                    dist = s_left - ref_right

        elif direction == "left":
            # 1. Est-ce à gauche ? (Droite du candidat <= Gauche de la ref)
            if s_right <= (ref_left + 5000):
                # 2. Est-ce aligné verticalement ?
                if s_top <= ref_center_y <= s_bottom:
                    match_found = True
                    dist = ref_left - s_right

        if match_found:
            candidates.append((dist, shape))

    # Tri par distance croissante (le plus proche d'abord)
    candidates.sort(key=lambda x: x[0])

    # Retourne le texte des N premiers candidats
    # Note: On filtre si pas de texteframe pour éviter les erreurs
    results = []
    for _, shape in candidates:
        if mode == "text":
            if hasattr(shape, "text_frame") and shape.text_frame and shape.text.strip():
                results.append(shape.text.strip())
        elif mode == "shape" :
            results.append(shape)
        if max_count != 0 and len(results) >= max_count:
            break
            
    return results

# ...existing code...

def get_shapes_inside(ref_shape, all_shapes, mode="text"):
    """
    Finds shapes that are visually contained within the boundaries of ref_shape.
    """
    if not ref_shape:
        return []

    # Reference boundaries
    ref_left = ref_shape.left
    ref_right = ref_shape.left + ref_shape.width
    ref_top = ref_shape.top
    ref_bottom = ref_shape.top + ref_shape.height
    
    results = []

    for shape in all_shapes:
        # Ignore self
        if shape.shape_id == ref_shape.shape_id:
            continue

        # Candidate boundaries
        s_left = shape.left
        s_right = shape.left + shape.width
        s_top = shape.top
        s_bottom = shape.top + shape.height

        # Start of strict containment check
        # A shape is inside if its Left is >= Ref Left AND Right <= Ref Right
        # AND Top >= Ref Top AND Bottom <= Ref Bottom
        
        # We add a tiny tolerance (margin) to handle border-perfect alignments
        margin = 0 # ~0.15mm
        
        is_inside_horizontally = (s_left >= ref_left - margin) and (s_right <= ref_right + margin)
        is_inside_vertically = (s_top >= ref_top - margin) and (s_bottom <= ref_bottom + margin)

        if is_inside_horizontally and is_inside_vertically:
             if mode == "text":
                if hasattr(shape, "text_frame") and shape.text_frame and shape.text.strip():
                    results.append(shape.text.strip())
             elif mode == "shape":
                results.append(shape)

    return results

def return_position(shape):
    return {
        "left": round(Length(shape.left).cm, 2),
        "right": round(Length(shape.left + shape.width).cm, 2),
        "top": round(Length(shape.top).cm, 2),
        "bottom": round(Length(shape.top + shape.height).cm, 2),
        "width": round(Length(shape.width).cm, 2),
        "height": round(Length(shape.height).cm, 2)
    }

def is_inside(ref_shape, target_shape, tolerance=0):
    """
    Vérifie si target_shape est visuellement à l'intérieur de ref_shape.
    
    Args:
        ref_shape: La forme conteneur (ex: le cadre d'une étape).
        target_shape: La forme à tester (ex: un texte).
        tolerance: Marge d'erreur en EMU (ex: 5000) pour accepter un léger dépassement.
    
    Returns:
        bool: True si target est dedans.
    """
    # 1. Calcul des limites de la Référence
    ref_left = ref_shape.left - tolerance
    ref_top = ref_shape.top - tolerance
    ref_right = ref_shape.left + ref_shape.width + tolerance
    ref_bottom = ref_shape.top + ref_shape.height + tolerance

    # 2. Calcul des limites de la Cible
    # Note: On doit gérer le cas où target est un Groupe (il a width/height) 
    # ou une forme simple.
    tgt_left = target_shape.left
    tgt_top = target_shape.top
    tgt_right = tgt_left + target_shape.width
    tgt_bottom = tgt_top + target_shape.height

    # 3. Vérification Géométrique
    # Est-ce que le coté gauche de la cible est APRES le coté gauche de la ref ?
    inside_horizontal = (tgt_left >= ref_left) and (tgt_right <= ref_right)
    
    # Est-ce que le haut de la cible est EN DESSOUS du haut de la ref ?
    inside_vertical = (tgt_top >= ref_top) and (tgt_bottom <= ref_bottom)

    return inside_horizontal and inside_vertical

def is_below(ref_shape, target_shape, tolerance=0):
    """
    Vérifie si target_shape est situé EN DESSOUS de ref_shape,
    tout en restant dans ses limites horizontales (effet colonne).
    
    Args:
        ref_shape: La forme du haut (ex: le Titre de l'étape).
        target_shape: La forme du bas (ex: le Groupe d'actions).
        tolerance: Marge en EMU pour l'alignement horizontal.
    """
    # 1. Limites de la colonne (basées sur la référence du haut)
    col_left = ref_shape.left - tolerance
    col_right = ref_shape.left + ref_shape.width + tolerance
    ref_bottom = ref_shape.top + ref_shape.height
    
    # 2. Position de la cible
    tgt_center_x = target_shape.left + (target_shape.width / 2)
    tgt_top = target_shape.top
    
    # 3. Vérification
    # A. Est-ce que c'est bien en dessous ? (Le haut de la cible >= Le bas de la ref)
    is_visually_below = tgt_top >= (ref_bottom - 5000) # Petite marge de 5000 EMU pour frôler
    
    # B. Est-ce que c'est aligné horizontalement ? (Le centre de la cible est entre gauche et droite de la ref)
    is_horizontally_aligned = (tgt_center_x >= col_left) and (tgt_center_x <= col_right)
    
    return is_visually_below and is_horizontally_aligned

def is_inside_range_horizontal(top, bot, target_shape, tolerance=0):
    """
    Vérifie si target_shape est aligné horizontalement dans la range top-bot,
    c'est-à-dire que son centre X est entre les bords gauche et droit de la zone définie par top et bot.
    
    Args:
        top: La position du top (ex: le Titre de l'étape).
        bot: La position du bot (ex: le Titre de l'étape suivante).
        target_shape: La forme à tester (ex: une action).
        tolerance: Marge en EMU pour élargir la zone d'alignement.
    """
    
    inside_vertical = (target_shape.top >= top - tolerance) and (target_shape.top + target_shape.height <= bot + tolerance)

    return inside_vertical