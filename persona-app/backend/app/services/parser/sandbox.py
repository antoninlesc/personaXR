import json
import os
from pptx.util import Length
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.connector import Connector
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from utils import collect_all_shapes, get_shapes_in_direction, get_shapes_inside

def has_bbox(shape: BaseShape) -> bool:
    return all(hasattr(shape, a) for a in ("left", "top", "width", "height"))

def right_edge(shape: BaseShape) -> int:
    return int(shape.left + shape.width)

def bottom_edge(shape: BaseShape) -> int:
    return int(shape.top + shape.height)

def vertical_overlap(a: BaseShape, b: BaseShape) -> bool:
    # vrai si les projections verticales se chevauchent
    return not (bottom_edge(a) <= int(b.top) or bottom_edge(b) <= int(a.top))

def horizontal_overlap(a: BaseShape, b: BaseShape) -> bool:
    # vrai si les projections horizontales se chevauchent
    return not (right_edge(a) <= int(b.left) or right_edge(b) <= int(a.left))

def shapes_below(slide, ref_shape: BaseShape, margin_emu: int = 0, require_overlap=True):
    """
    Retourne les shapes dont le bord supérieur est sous le bord inférieur de ref_shape.
    margin_emu : marge (en EMU) pour éviter les cas "collés".
    require_overlap : si True, on garde seulement celles qui chevauchent horizontalement ref_shape.
    """
    ref_bottom = bottom_edge(ref_shape) + margin_emu

    candidates = []
    for sh in slide.shapes:
        if sh is ref_shape:
            continue
        if not has_bbox(sh) or not has_bbox(ref_shape):
            continue

        if int(sh.top) >= ref_bottom:
            if not require_overlap or horizontal_overlap(ref_shape, sh):
                candidates.append(sh)

    # tri: d'abord la distance verticale, puis la position horizontale
    candidates.sort(key=lambda s: (int(s.top) - ref_bottom, int(s.left)))
    return candidates

def shapes_to_the_right(slide, ref_shape: BaseShape, margin_emu: int = 0, require_overlap=True):
    """
    Retourne les shapes dont le bord gauche est à droite du bord droit de ref_shape.
    margin_emu : marge (en EMU) pour éviter les cas "collés".
    require_overlap : si True, on garde seulement celles qui chevauchent verticalement ref_shape.
    """
    ref_right = right_edge(ref_shape) + margin_emu

    candidates = []
    for sh in slide.shapes:
        if sh is ref_shape:
            continue
        if not has_bbox(sh) or not has_bbox(ref_shape):
            continue

        if int(sh.left) >= ref_right:
            if not require_overlap or vertical_overlap(ref_shape, sh):
                candidates.append(sh)

    # souvent utile : trier par distance puis par hauteur
    candidates.sort(key=lambda s: (int(s.left) - ref_right, int(s.top)))
    return candidates



def extract_pptx_structured(file_path):
    presentation = Presentation(file_path)
    slides_data = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_data = {
            "slide_number": slide_index,
            "title": None,
            "text_blocks": [],
            "right": [],
            "tables": [],
            "notes": None
        }

        # Extract title
        if slide.shapes.title:
            slide_data["title"] = slide.shapes.title.text.strip()

        for shape in slide.shapes:
            
            # Regular text
            if hasattr(shape, "text") and shape.text.strip():
                # Avoid duplicating title
                if shape != slide.shapes.title:
                    slide_data["text_blocks"].append(shape.text.strip())

            # Tables
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                slide_data["tables"].append(table_data)

            if shape.shape_id == 4:
                for sh in shapes_to_the_right(slide, shape, margin_emu=0, require_overlap=True):
                    slide_data["right"].append(sh.text.strip())

        # Notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_data["notes"] = notes_text

        
        slides_data.append(slide_data)

    return slides_data

def save_json(data, pptx_path, output_path=None):
    # Default: same folder as PPTX, same name but .json
    if output_path is None:
        base, _ = os.path.splitext(pptx_path)
        output_path = base + "sandbox" + ".json"

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)

    return output_path

def collect_all_shapes(shapes):
    """Recursively collect all shapes, including those in groups"""
    all_shapes = []
    for shape in shapes:
        if shape.shape_type == 6: # Group
            all_shapes.extend(collect_all_shapes(shape.shapes))
        else:
            all_shapes.append(shape)
    return all_shapes

def get_shapes_below(title_shape, all_shapes, max_count=1):
    if not title_shape:
        return []

    # 1. Le "Laser" part du centre du titre
    title_center_x = title_shape.left + (title_shape.width / 2)
    title_bottom = title_shape.top + title_shape.height

    candidates = []

    for shape in all_shapes:
        # On ignore le titre lui-même
        if shape.shape_id == title_shape.shape_id:
            continue
            
        # 2. Vérification Verticale : La forme doit être EN DESSOUS du titre
        # On ajoute une petite tolérance (5mm = ~18000 EMU) au cas où les blocs se chevauchent à peine
        if shape.top >= (title_bottom - 18000):
            
            # 3. Vérification Horizontale : Le "Laser" traverse-t-il la forme ?
            # Est-ce que le centre du titre est COMPRIS entre le bord gauche et le bord droit de la forme ?
            shape_left = shape.left
            shape_right = shape.left + shape.width
            
            # Condition de croisement : Left < CenterX < Right
            if shape_left <= title_center_x <= shape_right:
                if hasattr(shape, "text") and shape.text.strip():
                # C'est un candidat valide ! On calcule la distance verticale
                    dist_y = shape.top - title_bottom
                    candidates.append((dist_y, shape))

    # 4. Tri : On prend les formes du haut vers le bas (les plus proches du titre d'abord)
    candidates.sort(key=lambda x: x[0])

    # On retourne le texte des N premiers candidats
    if max_count == 0:
        return [c[1] for c in candidates
                ]
    return [c[1] for c in candidates[:max_count]]

def get_all_shapes_by_group(slide):
    """Retourne un dict avec les formes groupées par leur parent (ou None si pas de groupe)"""
    shape_dict = {"groups": []}
    for shape in slide.shapes:
        if shape.shape_type == 6: # Group
            shape_dict["groups"].append({shape.name: [f"Left pos: {Length(shape.left).cm:.2f}", f"Right pos: {Length(shape.left + shape.width).cm:.2f}", f"Top pos: {Length(shape.top).cm:.2f}", f"Width: {Length(shape.width).cm:.2f}", f"Height: {Length(shape.height).cm:.2f}"  ]})
            for sub in shape.shapes:
                if shape.name not in shape_dict:
                    shape_dict[shape.name] = []
                shape_dict[shape.name].append(sub.name)
        else:
            shape_dict[shape.name] = None
    return shape_dict

if __name__ == "__main__":
    file_path = r"C:\Users\Laffineur\Documents\school\m2\ProjectPersona\pptx\pptx_example_cindy.pptx"
    pres = Presentation(file_path)
    slide = pres.slides[1]

    all_shapes = collect_all_shapes(slide.shapes)
    all_shapes_by_group = get_all_shapes_by_group(slide)
    print(f"Total shapes (including groups): {len(all_shapes)}")
    save_json(all_shapes_by_group, file_path)
    # print(f"Total shapes (including groups): {len(all_shapes)}")
    # print("All shapes:")
    every_shapes = []
    for shapes in slide.shapes:
        every_shapes.append(shapes)

    data = []

    for shape in all_shapes:
        # print(f" {shape.name} (type: {shape.shape_type})")
        # print(f"     Text: {getattr(shape, 'text', '')}")
        # if shape.name == "Rectangle 172":
        #     print(f"Found '{shape.name}' with text: '{getattr(shape, 'text', '')}'")
        #     print(f"Shape details - Left: {Length(shape.left).cm:.2f} cm, Top: {Length(shape.top).cm:.2f} cm, Width: {Length(shape.width).cm:.2f} cm, Height: {Length(shape.height).cm:.2f} cm")
        #     below_shapes = get_shapes_in_direction(shape, all_shapes, max_count=0, direction="below", mode="shape")
        #     print(f"Shapes below '{shape.name}': {[s.name for s in below_shapes]}")
        #     print(f"Text of shapes below '{shape.name}': {[getattr(s, 'text', '') for s in below_shapes]}")
            
        #     print("================================")

        if hasattr(shape, "text_frame") and shape.text_frame:
            if "etape" in shape.text_frame.text.lower():
                print(f" {shape.name} (type: {shape.shape_type})     --> Contains 'étape'")
                text = get_shapes_in_direction(shape, all_shapes, max_count=0, direction="below", mode="shape")
                print(f"Shapes below '{shape.name}': {[s.name for s in text]}")
                print(f"Text of shapes below '{shape.name}': {[getattr(s, 'text', '') for s in text]}")

                print("================================")

                etape_shape = shape
                action_shape = text[0] if len(text) > 0 else None
                canaux_shape = text[1] if len(text) > 1 else None
                emotions_shape = text[2] if len(text) > 2 else None
                challenges_shape = text[3] if len(text) > 3 else None

                # right_of_actions = get_shapes_in_direction(action_shape, all_shapes, direction="right", max_count=0, mode="shape")
                # for s in right_of_actions:
                #     insides = get_shapes_inside(s, every_shapes, mode='shape')
                #     inside = {
                #         s.name: []
                #     }

                #     for ins in insides:
                #         if hasattr(ins, "shapes"):
                #             for sub in ins.shapes:
                #                 print(f"sub text: {getattr(sub, 'text', '').strip("")}")
                #                 if hasattr(sub, "text") and sub.text.strip():
                #                     inside[s.name].append(sub.text.strip())

                #         else :
                #             inside[s.name].append(getattr(ins, "text", "").strip())
                #     data.append(inside)
                # save_json(data, file_path, output_path=r"C:\Users\Laffineur\Documents\school\m2\ProjectPersona\pptx\extracted_slide2.json")

                # print(f"Shapes to the right of '{action_shape.name}': {[s.name for s in right_of_actions]}")
                # print(f"Shapes inside '{right_of_actions[1].name}': {[s.name for s in insides]}")
                # print(f"Text of shapes inside '{right_of_actions[1].name}': {[getattr(s, 'text', '') for s in insides]}")

                
    #             print("================================")
    # # right_of_etape =  get_shapes_in_direction(etape_shape, all_shapes, direction="right", max_count=0, mode="shape")
    # # print("ETAPES : ")
    # # print(f"Shapes to the right of '{etape_shape.name}': {[s.name for s in right_of_etape]}")
    # # print("===================================")

    # right_of_actions = get_shapes_in_direction(action_shape, all_shapes, direction="right", max_count=0, mode="shape")
    # print("ACTIONS : ")
    # print(f"Shapes to the right of '{action_shape.name}': {[s.name for s in right_of_actions]}")
    
    # print(f" Text inside '{right_of_actions[0].name}': ")
    # insides = get_shapes_inside(right_of_actions[0], all_shapes, mode='shape')
    # print(f"Shapes inside {right_of_actions[0].name}")
    # for s in insides:
    #     print(f" - {s.name} (type: {s.shape_type}) with text: '{getattr(s, 'text', '')}'")
    
    print("===================================")

    right_of_challenges = get_shapes_in_direction(challenges_shape, all_shapes, direction="right", max_count=0)
    print("CHALLENGES : ")
    print(f"Shapes to the right of '{challenges_shape.name}': {[s for s in right_of_challenges]}")
    print("===================================")
            # if "bio" in shape.text_frame.text.lower():
            #     print(f" {shape.name} (type: {shape.shape_type})     --> Contains 'bio'")
    
            #     text = get_shapes_below(shape, all_shapes)
            #     print(f"Shapes below '{shape.name}': {[s.name for s in text]}")
            #     print(f"Text of shapes below '{shape.name}': {[getattr(s, 'text', '') for s in text]}")

            #     print("================================")
            
            # if "inquiétudes" in shape.text_frame.text.lower():
            #     print(f" {shape.name} (type: {shape.shape_type})     --> Contains 'inquiétudes'")
            #     text = get_shapes_below(shape, all_shapes)
            #     print(f"Shapes below '{shape.name}': {[s.name for s in text]}")
            #     print(f"Shapes below has text: {[hasattr(s, 'text') for s in text]}")
            #     print(f"Text of shapes below '{shape.name}': {[getattr(s, 'text', '') for s in text]}")
            
            #     print("================================")
            
            # if "besoins" in shape.text_frame.text.lower():
            #     print(f" {shape.name} (type: {shape.shape_type})     --> Contains 'besoins'")
            #     text = get_shapes_below(shape, all_shapes)
            #     print(f"Shapes below '{shape.name}': {[s.name for s in text]}")
            #     print(f"Text of shapes below '{shape.name}': {[getattr(s, 'text', '') for s in text]}")

            #     print("================================")
            
            # if "insight" in shape.text_frame.text.lower():
            #     print(f" {shape.name} (type: {shape.shape_type})     --> Contains 'insight'")
            #     text = get_shapes_below(shape, all_shapes, 2)
            #     print(f"Shapes below '{shape.name}': {[s.name for s in text]}")
            #     print(f"Text of shapes below '{shape.name}': {[getattr(s, 'text', '') for s in text]}")
            #     print("================================")
            
    # structured_output = extract_pptx_structured(file_path)
    # save_json(structured_output, file_path)
